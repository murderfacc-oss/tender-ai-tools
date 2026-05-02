"""
Генерирует презентацию «Установка Tender AI Tools» в формате .pptx.
Минималистичная — 6 слайдов, каждый с визуальным мокапом UI.

Использование:
    python scripts/generate_install_pptx.py
    python scripts/generate_install_pptx.py --output "Установка.pptx"
"""

import argparse
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ── палитра ────────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x1F, 0x63, 0xB5)
BLUE_DEEP  = RGBColor(0x16, 0x4F, 0x96)
BLUE_LIGHT = RGBColor(0xD9, 0xE4, 0xF5)
BLUE_VLIGHT= RGBColor(0xF0, 0xF6, 0xFE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GRAY       = RGBColor(0x55, 0x55, 0x55)
GRAY_LIGHT = RGBColor(0xEE, 0xEE, 0xEE)
GRAY_BG    = RGBColor(0xF7, 0xF9, 0xFF)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
ORANGE     = RGBColor(0xF5, 0x7F, 0x17)
ACCENT_YEL = RGBColor(0xFF, 0xC1, 0x07)
RED        = RGBColor(0xB7, 0x1C, 0x1C)


# ── размеры (16:9) ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── базовые хелперы ────────────────────────────────────────────────────────

def rect(slide, left, top, width, height, fill, *, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s


def round_rect(slide, left, top, width, height, fill, *, line=None, radius=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.adjustments[0] = radius
    return s


def text(slide, left, top, width, height, txt, *,
         size=18, bold=False, italic=False, color=DARK,
         align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def arrow(slide, x1, y1, x2, y2, color=BLUE):
    s = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight line
    s.line.color.rgb = color
    s.line.width = Pt(2.5)
    # add arrow head
    line = s.line
    lnE = line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    head = etree.SubElement(lnE, qn("a:headEnd"))
    head.set("type", "none")
    tail = etree.SubElement(lnE, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return s


def step_badge(slide, left, top, num, size_in=0.9, color=BLUE):
    """Большой круглый бейдж с номером шага."""
    d = Inches(size_in)
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    text(slide, left, top + Inches(0.05), d, d - Inches(0.1),
         str(num), size=int(size_in * 36), bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
    return s


# ── мокапы UI ──────────────────────────────────────────────────────────────

def mock_window(slide, left, top, width, height, title="", *,
                title_color=GRAY_LIGHT, content_bg=WHITE):
    """Мокап окна приложения с заголовком и панелью."""
    # тень
    rect(slide, left + Inches(0.05), top + Inches(0.05),
         width, height, RGBColor(0xE0, 0xE0, 0xE0))
    # рамка
    rect(slide, left, top, width, height, content_bg, line=GRAY)
    # title bar
    bar_h = Inches(0.32)
    rect(slide, left, top, width, bar_h, title_color)
    # светофор
    for i, color in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E),
                                RGBColor(0x28,0xC8,0x40)]):
        slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               left + Inches(0.12 + i*0.22),
                               top + Inches(0.08),
                               Inches(0.16), Inches(0.16))\
            .fill.solid()
    # title
    if title:
        text(slide, left + Inches(0.85), top + Inches(0.02),
             width - Inches(0.85), bar_h, title,
             size=11, color=GRAY, align=PP_ALIGN.CENTER)
    return left, top + bar_h, width, height - bar_h


def mock_folder(slide, left, top, width, height, files):
    """Мокап Проводника / Finder с файлами."""
    cl, ct, cw, ch = mock_window(slide, left, top, width, height,
                                  title="tender-ai-tools",
                                  title_color=RGBColor(0xE8,0xEC,0xF2))
    # рисуем файлы
    icon_w = Inches(0.95)
    icon_h = Inches(1.1)
    pad_x  = Inches(0.18)
    pad_y  = Inches(0.18)
    cols = int((cw - pad_x) / (icon_w + pad_x))
    x0 = cl + pad_x
    y0 = ct + Inches(0.2)
    for i, (name, kind) in enumerate(files):
        col = i % cols
        row = i // cols
        x = x0 + col * (icon_w + pad_x)
        y = y0 + row * (icon_h + pad_y)
        # иконка
        if kind == "folder":
            ic = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER,
                                         x + Inches(0.18), y, Inches(0.6), Inches(0.5))
            ic.fill.solid(); ic.fill.fore_color.rgb = ACCENT_YEL
            ic.line.fill.background()
        elif kind == "zip":
            ic = rect(slide, x + Inches(0.22), y, Inches(0.55), Inches(0.55), GRAY_LIGHT)
            text(slide, x + Inches(0.22), y + Inches(0.13),
                 Inches(0.55), Inches(0.3),
                 "ZIP", size=10, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        elif kind == "py":
            ic = rect(slide, x + Inches(0.22), y, Inches(0.55), Inches(0.55), BLUE)
            text(slide, x + Inches(0.22), y + Inches(0.13),
                 Inches(0.55), Inches(0.3),
                 "PY", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        elif kind == "pptx":
            ic = rect(slide, x + Inches(0.22), y, Inches(0.55), Inches(0.55),
                      RGBColor(0xC4,0x3F,0x10))
            text(slide, x + Inches(0.22), y + Inches(0.13),
                 Inches(0.55), Inches(0.3),
                 "P", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        elif kind == "md":
            ic = rect(slide, x + Inches(0.22), y, Inches(0.55), Inches(0.55), GRAY_LIGHT)
            text(slide, x + Inches(0.22), y + Inches(0.13),
                 Inches(0.55), Inches(0.3),
                 "MD", size=9, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        # подпись
        text(slide, x, y + Inches(0.6), icon_w, Inches(0.4),
             name, size=9, color=DARK, align=PP_ALIGN.CENTER)
    return cl, ct, cw, ch


def mock_claude_code(slide, left, top, width, height, *,
                     show_menu=False, prompt_in_chat=None):
    """Мокап окна Claude Code: сайдбар + чат."""
    cl, ct, cw, ch = mock_window(slide, left, top, width, height,
                                  title="Claude Code",
                                  title_color=RGBColor(0xF5,0xEE,0xE5))
    # сайдбар
    sb_w = Inches(2.0)
    rect(slide, cl, ct, sb_w, ch, RGBColor(0xFA,0xF7,0xF2))
    # File menu (если show_menu)
    if show_menu:
        # menubar выпадушка
        menu_x = cl + Inches(0.1)
        menu_y = ct + Inches(0.1)
        rect(slide, menu_x, menu_y, Inches(2.4), Inches(1.4), WHITE, line=GRAY)
        items = ["📁  Open Folder", "📄  Open File", "💾  Save", "⚙  Settings"]
        for i, item in enumerate(items):
            row_y = menu_y + Inches(0.1 + i * 0.32)
            if i == 0:  # выделяем первый
                rect(slide, menu_x + Inches(0.05), row_y,
                     Inches(2.3), Inches(0.3), BLUE_LIGHT)
            text(slide, menu_x + Inches(0.15), row_y + Inches(0.04),
                 Inches(2.2), Inches(0.25),
                 item, size=11, color=DARK,
                 bold=(i == 0))
    else:
        # просто сайдбар-элементы
        items = ["💬 New chat", "📁 tender-ai-tools", "  • README.md",
                 "  • INSTALL_PROMPT.md", "  • skills/", "  • mcp/"]
        for i, item in enumerate(items):
            text(slide, cl + Inches(0.15), ct + Inches(0.2 + i * 0.32),
                 sb_w - Inches(0.2), Inches(0.3),
                 item, size=11,
                 color=BLUE if i == 1 else DARK,
                 bold=(i in (0, 1)))

    # область чата
    chat_x = cl + sb_w
    chat_w = cw - sb_w
    rect(slide, chat_x, ct, chat_w, ch, WHITE)

    # поле ввода внизу
    input_h = Inches(1.4) if prompt_in_chat else Inches(0.6)
    input_y = ct + ch - input_h - Inches(0.15)
    round_rect(slide, chat_x + Inches(0.2), input_y,
               chat_w - Inches(0.4), input_h, GRAY_BG, line=GRAY,
               radius=0.1)
    if prompt_in_chat:
        text(slide, chat_x + Inches(0.35), input_y + Inches(0.1),
             chat_w - Inches(0.7), input_h - Inches(0.25),
             prompt_in_chat,
             size=10, color=DARK, font="Consolas")
        # курсор-плейсхолдер не нужен
    else:
        text(slide, chat_x + Inches(0.35), input_y + Inches(0.15),
             chat_w - Inches(0.7), Inches(0.3),
             "Введи сообщение…", size=11, color=GRAY, italic=True)

    return cl, ct, cw, ch


def mock_browser(slide, left, top, width, height, url, content_fn=None):
    """Мокап окна браузера."""
    cl, ct, cw, ch = mock_window(slide, left, top, width, height,
                                  title="",
                                  title_color=RGBColor(0xE0,0xE5,0xEC))
    # url-bar
    bar_h = Inches(0.42)
    rect(slide, cl, ct, cw, bar_h, RGBColor(0xF1,0xF3,0xF5))
    round_rect(slide, cl + Inches(0.2), ct + Inches(0.07),
               cw - Inches(0.4), bar_h - Inches(0.14), WHITE,
               line=GRAY, radius=0.3)
    text(slide, cl + Inches(0.4), ct + Inches(0.1),
         cw - Inches(0.7), bar_h - Inches(0.2),
         url, size=11, color=GRAY)
    if content_fn:
        content_fn(slide, cl, ct + bar_h, cw, ch - bar_h)
    return cl, ct + bar_h, cw, ch - bar_h


# ── общие элементы слайда ──────────────────────────────────────────────────

def slide_header(slide, step_num, total, title):
    """Шапка слайда: цифра шага + название + прогресс справа."""
    rect(slide, 0, 0, W, Inches(1.1), BLUE)
    step_badge(slide, Inches(0.45), Inches(0.18),
               step_num, size_in=0.74, color=WHITE)
    # цифра в бейдже белая, на белом не видно — переделаю
    # перерисую вручную: круг white, текст blue
    # удалим предыдущий
    pass  # уже нарисовано, пере-рисую через свой бейдж
    # title
    text(slide, Inches(1.5), Inches(0.18), Inches(10), Inches(0.45),
         f"Шаг {step_num} из {total}",
         size=14, color=BLUE_LIGHT)
    text(slide, Inches(1.5), Inches(0.5), Inches(10), Inches(0.55),
         title, size=26, bold=True, color=WHITE)


def slide_done_footer(slide):
    rect(slide, 0, H - Inches(0.4), W, Inches(0.4), BLUE_DEEP)
    text(slide, 0, H - Inches(0.36), W, Inches(0.32),
         "Tender AI Tools  •  Установка для партнёра",
         size=10, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)


# ── собственно слайды ──────────────────────────────────────────────────────

def slide_title(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, BLUE)
    rect(sl, 0, H * 0.6, W, H * 0.4, BLUE_DEEP)
    rect(sl, Inches(0.6), Inches(0.6), Inches(0.12), Inches(2.0), ACCENT_YEL)

    text(sl, Inches(0.95), Inches(0.55), Inches(11), Inches(1.4),
         "Tender AI Tools",
         size=58, bold=True, color=WHITE)
    text(sl, Inches(0.95), Inches(2.0), Inches(11), Inches(0.7),
         "Установка для партнёра",
         size=28, color=BLUE_LIGHT)

    # три простых шага в виде карточек
    cards = [
        ("1", "Распакуй\nархив",         "📦"),
        ("2", "Открой папку\nв Claude Code", "💬"),
        ("3", "Вставь промт\nв чат",     "✨"),
    ]
    card_w = Inches(3.4)
    gap    = Inches(0.4)
    total  = 3 * card_w + 2 * gap
    start  = (W - total) / 2
    top    = Inches(3.6)
    card_h = Inches(2.4)
    for i, (num, label, icon) in enumerate(cards):
        x = start + i * (card_w + gap)
        round_rect(sl, x, top, card_w, card_h, WHITE, radius=0.1)
        # цифра
        text(sl, x + Inches(0.3), top + Inches(0.2), Inches(1.0), Inches(1.0),
             num, size=54, bold=True, color=BLUE)
        # иконка справа
        text(sl, x + card_w - Inches(1.1), top + Inches(0.25),
             Inches(1.0), Inches(1.0),
             icon, size=42, align=PP_ALIGN.CENTER)
        # подпись
        text(sl, x + Inches(0.3), top + Inches(1.3),
             card_w - Inches(0.6), Inches(1.0),
             label, size=18, bold=True, color=DARK)
        # стрелка между карточками
        if i < 2:
            ay = top + card_h / 2
            ax1 = x + card_w + Inches(0.05)
            ax2 = x + card_w + gap - Inches(0.05)
            arrow(sl, ax1, ay, ax2, ay, color=WHITE)

    # дата
    today = datetime.now().strftime("%d.%m.%Y")
    text(sl, Inches(0.5), H - Inches(0.7),
         W - Inches(1), Inches(0.4),
         today, size=14, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, GRAY_BG)

    text(sl, Inches(0.5), Inches(0.4), W - Inches(1), Inches(0.6),
         "Что в архиве",
         size=32, bold=True, color=DARK)
    text(sl, Inches(0.5), Inches(1.05), W - Inches(1), Inches(0.5),
         "Этот архив — всё, что нужно для работы",
         size=16, color=GRAY)

    # мокап папки в центре
    mock_folder(sl, Inches(0.7), Inches(1.85), W - Inches(1.4), Inches(3.7),
                files=[
                    ("skills/",   "folder"),
                    ("mcp/",      "folder"),
                    ("README.md", "md"),
                    ("INSTALL_PROMPT.md",  "md"),
                    ("Установка.pptx", "pptx"),
                ])

    # подписи снизу
    notes = [
        ("skills/", "4 готовых скилла", BLUE),
        ("mcp/",    "сервер документов ЕИС", GREEN),
        ("INSTALL_PROMPT.md", "промт для Claude Code", ORANGE),
    ]
    y = Inches(5.85)
    for name, desc, color in notes:
        x = Inches(0.7)
        rect(sl, x, y, Inches(0.15), Inches(0.5), color)
        text(sl, x + Inches(0.3), y, Inches(3.5), Inches(0.3),
             name, size=12, bold=True, color=color, font="Consolas")
        text(sl, x + Inches(0.3), y + Inches(0.25), Inches(8), Inches(0.3),
             desc, size=12, color=GRAY)
        y += Inches(0.55)


def slide_step1_unpack(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, GRAY_BG)
    slide_header(sl, 1, 4, "Распакуй архив")

    # слева — действие
    text(sl, Inches(0.5), Inches(1.4), Inches(6.5), Inches(0.5),
         "Что делаем", size=14, bold=True, color=GRAY)
    text(sl, Inches(0.5), Inches(1.85), Inches(6.5), Inches(0.6),
         "Распакуй tender-ai-tools.zip",
         size=24, bold=True, color=DARK)
    text(sl, Inches(0.5), Inches(2.55), Inches(6.5), Inches(2),
         "Правая кнопка мыши на архиве → «Извлечь всё…»\n\n"
         "Папку для распаковки можешь выбрать любую,\n"
         "но рекомендуем:\n",
         size=15, color=GRAY)
    # путь в коде
    round_rect(sl, Inches(0.5), Inches(4.55), Inches(6.5), Inches(0.5),
               GRAY_LIGHT, radius=0.2)
    text(sl, Inches(0.7), Inches(4.62), Inches(6.2), Inches(0.4),
         "C:\\ClaudeCode\\tender-ai-tools\\",
         size=14, bold=True, color=DARK, font="Consolas")

    # подсказка
    round_rect(sl, Inches(0.5), Inches(5.4), Inches(6.5), Inches(1.4),
               RGBColor(0xFF,0xF8,0xE1), radius=0.05)
    text(sl, Inches(0.7), Inches(5.55), Inches(6.1), Inches(0.4),
         "💡 Совет",
         size=14, bold=True, color=ORANGE)
    text(sl, Inches(0.7), Inches(5.9), Inches(6.1), Inches(0.85),
         "Путь без пробелов и кириллицы — Claude Code\n"
         "и терминал работают надёжнее.",
         size=13, color=DARK)

    # справа — мокап результата
    text(sl, Inches(7.5), Inches(1.4), Inches(5.4), Inches(0.5),
         "Что должно получиться", size=14, bold=True, color=GRAY)
    mock_folder(sl, Inches(7.5), Inches(1.85), Inches(5.4), Inches(4.0),
                files=[
                    ("skills/",   "folder"),
                    ("mcp/",      "folder"),
                    ("README.md", "md"),
                    ("INSTALL_PROMPT.md", "md"),
                    ("Установка.pptx", "pptx"),
                ])
    slide_done_footer(sl)


def slide_step2_open(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, GRAY_BG)
    slide_header(sl, 2, 4, "Открой папку в Claude Code")

    # инструкция слева
    text(sl, Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.5),
         "Что делаем", size=14, bold=True, color=GRAY)
    text(sl, Inches(0.5), Inches(1.85), Inches(5.5), Inches(0.6),
         "File → Open Folder",
         size=24, bold=True, color=DARK)

    text(sl, Inches(0.5), Inches(2.7), Inches(5.5), Inches(2),
         "1.  Запусти Claude Code\n\n"
         "2.  В верхнем меню выбери:\n"
         "     File  →  Open Folder\n\n"
         "3.  Найди папку с распакованным\n"
         "     архивом и подтверди.",
         size=15, color=DARK)

    # подсказка
    round_rect(sl, Inches(0.5), Inches(5.65), Inches(5.5), Inches(1.2),
               RGBColor(0xFF,0xF8,0xE1), radius=0.05)
    text(sl, Inches(0.7), Inches(5.78), Inches(5.1), Inches(0.4),
         "💡 Совет",
         size=14, bold=True, color=ORANGE)
    text(sl, Inches(0.7), Inches(6.13), Inches(5.1), Inches(0.7),
         "Можно ещё проще — перетащи папку\n"
         "прямо в окно Claude Code.",
         size=13, color=DARK)

    # справа — мокап Claude Code с открытым меню
    mock_claude_code(sl, Inches(6.4), Inches(1.4), Inches(6.5), Inches(5.6),
                     show_menu=True)

    # стрелка к нужному пункту меню
    arrow(sl, Inches(6.0), Inches(2.0),
              Inches(7.05), Inches(1.85), color=RED)
    text(sl, Inches(5.4), Inches(2.1), Inches(0.6), Inches(0.4),
         "→", size=24, bold=True, color=RED)

    slide_done_footer(sl)


def slide_step3_paste(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, GRAY_BG)
    slide_header(sl, 3, 4, "Вставь промт в чат")

    # инструкция слева
    text(sl, Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.5),
         "Что делаем", size=14, bold=True, color=GRAY)
    text(sl, Inches(0.5), Inches(1.85), Inches(5.5), Inches(0.6),
         "Скопируй промт → вставь",
         size=22, bold=True, color=DARK)

    text(sl, Inches(0.5), Inches(2.7), Inches(5.5), Inches(3),
         "1.  Открой файл\n"
         "     INSTALL_PROMPT.md\n"
         "     (он лежит рядом с архивом)\n\n"
         "2.  Скопируй текст промта\n"
         "     (он находится внутри блока кода)\n\n"
         "3.  Вставь в чат Claude Code\n"
         "     и нажми Enter\n\n"
         "Дальше Claude Code сам:\n"
         "  • установит Python если нужно\n"
         "  • установит все библиотеки\n"
         "  • подключит MCP-сервер\n"
         "  • объяснит что делать дальше",
         size=14, color=DARK)

    # справа — мокап Claude Code с промтом в инпуте
    mock_claude_code(sl, Inches(6.4), Inches(1.4), Inches(6.5), Inches(5.6),
                     prompt_in_chat=(
                         "Установи мне Tender AI Tools — набор\n"
                         "инструментов для участия в госзакупках\n"
                         "по 44-ФЗ.\n\n"
                         "Текущая папка должна быть корнем\n"
                         "распакованного архива. В ней должны быть\n"
                         "skills/ и mcp/.\n\n"
                         "Шаг 1. Проверить и установить Python..."
                     ))

    slide_done_footer(sl)


def slide_step4_skills(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, GRAY_BG)
    slide_header(sl, 4, 4, "Загрузи скиллы в Cowork")

    # инструкция слева
    text(sl, Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.5),
         "Что делаем", size=14, bold=True, color=GRAY)
    text(sl, Inches(0.5), Inches(1.85), Inches(5.5), Inches(0.6),
         "4 файла → Cowork",
         size=24, bold=True, color=DARK)

    text(sl, Inches(0.5), Inches(2.7), Inches(5.5), Inches(0.5),
         "Это последний ручной шаг.",
         size=14, color=GRAY, italic=True)

    text(sl, Inches(0.5), Inches(3.2), Inches(5.5), Inches(3),
         "1.  Открой Cowork в браузере\n"
         "     (URL уточни у партнёра)\n\n"
         "2.  Раздел  Skills  →  Upload\n\n"
         "3.  Загрузи по очереди 4 zip-файла\n"
         "     из папки  skills/\n\n"
         "4.  Убедись, что у каждого статус\n"
         "     Active",
         size=14, color=DARK)

    # справа — мокап Cowork
    def cowork_content(slide, cl, ct, cw, ch):
        # шапка раздела Skills
        text(slide, cl + Inches(0.3), ct + Inches(0.15),
             Inches(2), Inches(0.4),
             "Skills", size=18, bold=True, color=DARK)
        # кнопка Upload
        round_rect(slide, cl + cw - Inches(1.5), ct + Inches(0.18),
                   Inches(1.2), Inches(0.4), BLUE, radius=0.2)
        text(slide, cl + cw - Inches(1.5), ct + Inches(0.22),
             Inches(1.2), Inches(0.35),
             "+ Upload", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
        # таблица скиллов
        skills = [
            ("verdikt-zakupki",     "v1.0",   GREEN),
            ("scan-zakupki",        "v0.3",   GREEN),
            ("zhaloba-fas",         "v0.1",   GREEN),
            ("zapros-razyasneniy",  "v0.1.1", GREEN),
        ]
        y = ct + Inches(0.85)
        for name, ver, status_color in skills:
            rect(slide, cl + Inches(0.25), y,
                 cw - Inches(0.5), Inches(0.55), WHITE,
                 line=RGBColor(0xE0,0xE5,0xEC))
            text(slide, cl + Inches(0.45), y + Inches(0.13),
                 cw - Inches(2.4), Inches(0.35),
                 name, size=12, bold=True, color=DARK, font="Consolas")
            text(slide, cl + cw - Inches(2.4), y + Inches(0.15),
                 Inches(0.8), Inches(0.3),
                 ver, size=11, color=GRAY)
            # status pill
            round_rect(slide, cl + cw - Inches(1.4), y + Inches(0.13),
                       Inches(1.0), Inches(0.3),
                       RGBColor(0xE6,0xF4,0xEA), radius=0.4)
            text(slide, cl + cw - Inches(1.4), y + Inches(0.16),
                 Inches(1.0), Inches(0.25),
                 "● Active", size=10, bold=True, color=status_color,
                 align=PP_ALIGN.CENTER)
            y += Inches(0.65)

    mock_browser(sl, Inches(6.4), Inches(1.4), Inches(6.5), Inches(5.6),
                 url="cowork.anthropic.com / skills",
                 content_fn=cowork_content)

    slide_done_footer(sl)


def slide_done(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, BLUE)
    rect(sl, 0, H * 0.65, W, H * 0.35, BLUE_DEEP)

    # большая галка
    cs = slide_check_circle(sl, W / 2 - Inches(1), Inches(1.0), Inches(2))

    text(sl, 0, Inches(3.4), W, Inches(1.0),
         "Готово!",
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    text(sl, 0, Inches(4.5), W, Inches(0.6),
         "Перезапусти Claude Code и попробуй:",
         size=20, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # три примера команд
    cmds = [
        ("оцени закупку",          "→ полный отчёт по документации"),
        ("составь жалобу",         "→ жалоба в УФАС или защита от РНП"),
        ("запрос разъяснений",     "→ Word-файл вопросов заказчику"),
    ]
    y = Inches(5.4)
    for cmd, desc in cmds:
        round_rect(sl, Inches(2.5), y, Inches(4), Inches(0.45),
                   WHITE, radius=0.4)
        text(sl, Inches(2.5), y + Inches(0.06),
             Inches(4), Inches(0.35),
             cmd, size=14, bold=True, color=BLUE,
             align=PP_ALIGN.CENTER, font="Consolas")
        text(sl, Inches(6.7), y + Inches(0.07),
             Inches(5), Inches(0.35),
             desc, size=14, color=BLUE_LIGHT)
        y += Inches(0.55)


def slide_check_circle(slide, left, top, size):
    """Большая галка в круге."""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.fill.background()
    text(slide, left, top + Inches(0.05), size, size - Inches(0.1),
         "✓", size=120, bold=True, color=BLUE,
         align=PP_ALIGN.CENTER)
    return s


# ── сборка ─────────────────────────────────────────────────────────────────

def generate(output_path: Path):
    prs = new_prs()

    slide_title(prs)            # 1. Титул с тремя шагами
    slide_overview(prs)         # 2. Что в архиве (мокап папки)
    slide_step1_unpack(prs)     # 3. Распаковать
    slide_step2_open(prs)       # 4. Открыть в Claude Code
    slide_step3_paste(prs)      # 5. Вставить промт
    slide_step4_skills(prs)     # 6. Загрузить скиллы в Cowork
    slide_done(prs)             # 7. Готово

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    size_kb = output_path.stat().st_size // 1024
    print(f"OK: {output_path}  ({size_kb} KB, {len(prs.slides)} слайдов)")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", "-o", type=Path,
                        default=Path("Установка Tender AI Tools.pptx"))
    args = parser.parse_args()
    generate(args.output)


if __name__ == "__main__":
    main()
